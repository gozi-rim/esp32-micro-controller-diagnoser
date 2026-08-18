"""
Generate Final_Project_Defence.pptx - 9-slide ECE 515.2 Project Defense Deck
for the NetDiag Expert System (ESP32 Hardware Diagnostics).

Design: 16:9 widescreen. Deep navy for ONLY slide 1 (title) and slide 9 (conclusion).
Content slides use clean white backgrounds with high-contrast dark text,
subtle muted accent borders, and gentle off-white card fills.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, datetime

# === PALETTE =================================================================
# Dark slides (1 and 9 only)
NAVY         = RGBColor(0x0A, 0x0C, 0x10)
DARK_CARD    = RGBColor(0x16, 0x1B, 0x22)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
MEDIUM_GRAY  = RGBColor(0x94, 0xA3, 0xB8)
MUTED_GRAY   = RGBColor(0x64, 0x74, 0x8B)

# Content slides (2-8) - clean white BG with strong dark text
SLIDE_BG     = RGBColor(0xFF, 0xFF, 0xFF)       # pure white
TITLE_TEXT   = RGBColor(0x11, 0x18, 0x27)        # near-black for headings
BODY_TEXT    = RGBColor(0x33, 0x40, 0x55)        # dark slate for body text
LEAD_TEXT    = RGBColor(0x1E, 0x29, 0x3B)        # dark for bold lead-ins
SUB_TEXT     = RGBColor(0x47, 0x55, 0x69)        # medium-dark for explanations

# Cards on white slides - very subtle tints (NOT saturated)
CARD_GRAY    = RGBColor(0xF8, 0xFA, 0xFC)        # barely-there gray
CARD_BORDER  = RGBColor(0xE2, 0xE8, 0xF0)        # light border

# Accent colors - used sparingly for tags, dots, thin bars
ACCENT_CYAN    = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_EMERALD = RGBColor(0x05, 0x9E, 0x6F)      # slightly muted emerald
ACCENT_BLUE    = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_AMBER   = RGBColor(0xD9, 0x77, 0x06)      # muted amber (not bright yellow)
ACCENT_PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
ACCENT_ROSE    = RGBColor(0xDB, 0x27, 0x77)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]


# === HELPERS =================================================================

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BODY_TEXT, align=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet(tf, bold_lead, explanation, font_size=16, lead_color=LEAD_TEXT,
               text_color=SUB_TEXT, dot_color=ACCENT_EMERALD, first=False,
               font_name="Calibri"):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(8)
    p.space_after  = Pt(4)
    p.line_spacing = Pt(font_size * 1.35)

    run_b = p.add_run()
    run_b.text = "   "
    run_b.font.size = Pt(10)
    run_b.font.color.rgb = dot_color
    run_b.font.name = font_name

    run_lead = p.add_run()
    run_lead.text = bold_lead + "  -  "
    run_lead.font.size = Pt(font_size)
    run_lead.font.bold = True
    run_lead.font.color.rgb = lead_color
    run_lead.font.name = font_name

    run_exp = p.add_run()
    run_exp.text = explanation
    run_exp.font.size = Pt(font_size)
    run_exp.font.bold = False
    run_exp.font.color.rgb = text_color
    run_exp.font.name = font_name
    return p


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.04
    return shape


def add_thin_bar(slide, color, y=Inches(7.25)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y, SLIDE_W, Inches(0.045)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_slide_number(slide, num, total=9, dark_bg=False):
    color = MUTED_GRAY if dark_bg else RGBColor(0xA0, 0xAE, 0xC0)
    add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
                f"{num} / {total}", font_size=11, color=color,
                align=PP_ALIGN.RIGHT)


def section_tag(slide, left, top, text, color=ACCENT_CYAN):
    add_textbox(slide, left, top, Inches(5), Inches(0.3),
                text.upper(), font_size=11, bold=True, color=color,
                font_name="Calibri")


def slide_title(slide, left, top, text, width=Inches(11)):
    add_textbox(slide, left, top, width, Inches(0.7),
                text, font_size=30, bold=True, color=TITLE_TEXT)


# =============================================================================
#  SLIDE 1 - TITLE  (Dark navy)
# =============================================================================
s1 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s1, NAVY)

# Left accent bar
add_rect(s1, Inches(0), Inches(0), Inches(0.1), SLIDE_H, ACCENT_CYAN, None)

# Tag
add_textbox(s1, Inches(1.2), Inches(0.9), Inches(6), Inches(0.3),
            "ECE 515.2  --  FORWARD-CHAINING EXPERT SYSTEM",
            font_size=12, bold=True, color=ACCENT_EMERALD)

# Main title
add_textbox(s1, Inches(1.2), Inches(1.9), Inches(11), Inches(1.4),
            "NetDiag Expert", font_size=52, bold=True, color=WHITE)

# Subtitle
add_textbox(s1, Inches(1.2), Inches(3.4), Inches(10), Inches(1.0),
            "A Localized Network & IoT Hardware Diagnostic Expert System\nfor ESP32-Based Embedded Systems",
            font_size=22, color=MEDIUM_GRAY, line_spacing=1.4)

# Divider
div = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.7), Inches(2.5), Inches(0.04))
div.fill.solid(); div.fill.fore_color.rgb = ACCENT_CYAN; div.line.fill.background()

# Presenter info
add_textbox(s1, Inches(1.2), Inches(5.1), Inches(7), Inches(0.4),
            "Presented by  --  Group 11", font_size=18, bold=True, color=WHITE)
add_textbox(s1, Inches(1.2), Inches(5.6), Inches(7), Inches(0.4),
            "Department of Electronic Engineering", font_size=15, color=MEDIUM_GRAY)
add_textbox(s1, Inches(1.2), Inches(6.0), Inches(7), Inches(0.4),
            "University of Port Harcourt", font_size=15, color=MEDIUM_GRAY)
add_textbox(s1, Inches(1.2), Inches(6.45), Inches(7), Inches(0.4),
            datetime.datetime.now().strftime("August %d, %Y"),
            font_size=14, color=MUTED_GRAY)

add_slide_number(s1, 1, dark_bg=True)


# =============================================================================
#  SLIDE 2 - PROBLEM STATEMENT  (White)
# =============================================================================
s2 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s2, SLIDE_BG)
add_thin_bar(s2, ACCENT_CYAN)

section_tag(s2, Inches(1.0), Inches(0.65), "Problem Statement")
slide_title(s2, Inches(1.0), Inches(1.1), "What Real-World Problem Does This Solve?")

# Card background
add_rect(s2, Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.7), CARD_GRAY, CARD_BORDER)

tf2 = add_rich_textbox(s2, Inches(1.5), Inches(2.5), Inches(10.3), Inches(4.3))

add_bullet(tf2,
    "Manual Debugging Is Slow & Error-Prone",
    "Engineers spend hours tracing ESP32 hardware faults using trial-and-error, often misdiagnosing power, GPIO, and communication issues.",
    first=True)
add_bullet(tf2,
    "No Structured Diagnostic Tools Exist",
    "Unlike software (which has debuggers and profilers), embedded hardware troubleshooting lacks accessible, intelligent, guided diagnostic platforms.")
add_bullet(tf2,
    "Expertise Gap in Educational Settings",
    "Students in ECE courses lack practical experience to quickly identify root causes of brownout resets, I2C lockups, or antenna interference.")
add_bullet(tf2,
    "Scattered Knowledge, No Central Reference",
    "Troubleshooting info is spread across datasheets, forums, and lecture notes with no single, searchable, decision-tree-driven knowledge base.")
add_bullet(tf2,
    "Real-Time Telemetry Is Inaccessible",
    "Monitoring live ESP32 signals typically requires expensive lab equipment that many students don't have access to.")

add_slide_number(s2, 2)


# =============================================================================
#  SLIDE 3 - PROJECT OVERVIEW & OBJECTIVES  (White)
# =============================================================================
s3 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s3, SLIDE_BG)
add_thin_bar(s3, ACCENT_EMERALD)

section_tag(s3, Inches(1.0), Inches(0.65), "Project Overview & Objectives", ACCENT_EMERALD)
slide_title(s3, Inches(1.0), Inches(1.1), "What We Built & The Main Goal")

# Left card
add_rect(s3, Inches(1.0), Inches(2.2), Inches(5.4), Inches(4.7), CARD_GRAY, CARD_BORDER)
add_textbox(s3, Inches(1.35), Inches(2.45), Inches(4.7), Inches(0.35),
            "PROJECT OVERVIEW", font_size=12, bold=True, color=ACCENT_EMERALD)

tf3l = add_rich_textbox(s3, Inches(1.35), Inches(3.0), Inches(4.7), Inches(3.6))
add_bullet(tf3l,
    "NetDiag Expert",
    "A web-based, rule-driven diagnostic expert system for ESP32 hardware troubleshooting.",
    font_size=15, first=True, dot_color=ACCENT_EMERALD)
add_bullet(tf3l,
    "Built With",
    "Next.js 16, React 19, TypeScript, Tailwind CSS - professional SaaS-style interface.",
    font_size=15, dot_color=ACCENT_EMERALD)
add_bullet(tf3l,
    "AI Methodology",
    "Forward-chaining inference engine walking a structured decision tree of 2,400+ diagnostic rules.",
    font_size=15, dot_color=ACCENT_EMERALD)

# Right card
add_rect(s3, Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.7), CARD_GRAY, CARD_BORDER)
add_textbox(s3, Inches(7.15), Inches(2.45), Inches(4.8), Inches(0.35),
            "KEY OBJECTIVES", font_size=12, bold=True, color=ACCENT_BLUE)

tf3r = add_rich_textbox(s3, Inches(7.15), Inches(3.0), Inches(4.8), Inches(3.6))
add_bullet(tf3r,
    "Automate Fault Isolation",
    "Replace manual trial-and-error with systematic, guided diagnostic workflows.",
    font_size=15, lead_color=LEAD_TEXT, first=True, dot_color=ACCENT_BLUE)
add_bullet(tf3r,
    "Provide Structured Knowledge",
    "Centralize 9 fault categories (Brownout, Wi-Fi, GPIO, I2C, SPI, ADC, etc.) into a searchable base.",
    font_size=15, lead_color=LEAD_TEXT, dot_color=ACCENT_BLUE)
add_bullet(tf3r,
    "Deliver Real-Time Monitoring",
    "Simulate live ESP32 telemetry (temperature, RSSI, heap) for educational insight.",
    font_size=15, lead_color=LEAD_TEXT, dot_color=ACCENT_BLUE)

add_slide_number(s3, 3)


# =============================================================================
#  SLIDE 4 - SYSTEM ARCHITECTURE  (White)
# =============================================================================
s4 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s4, SLIDE_BG)
add_thin_bar(s4, ACCENT_BLUE)

section_tag(s4, Inches(1.0), Inches(0.65), "System Architecture", ACCENT_BLUE)
slide_title(s4, Inches(1.0), Inches(1.1), "How It Works - High-Level Workflow")

steps = [
    ("STEP 1", "Describe the Symptom",
     "The user types or selects a hardware symptom (e.g., \"ESP32 keeps resetting\"). The system matches keywords to diagnostic entry points in the rule base.",
     ACCENT_CYAN),
    ("STEP 2", "Engine Processes Rules",
     "The forward-chaining inference engine fires matching rules, walks the decision tree, and narrows the fault domain step-by-step using IF-THEN logic.",
     ACCENT_EMERALD),
    ("STEP 3", "Receive Diagnosis",
     "The system outputs a structured diagnosis: root cause, confidence score, severity level, step-by-step remediation, and circuit diagram notes.",
     ACCENT_BLUE),
]

for i, (label, title, desc, accent) in enumerate(steps):
    left = Inches(1.0 + i * 3.85)
    add_rect(s4, left, Inches(2.2), Inches(3.5), Inches(4.7), CARD_GRAY, CARD_BORDER)

    # Colored top bar inside card
    bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.2), Inches(3.5), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    add_textbox(s4, Inches(left.inches + 0.3), Inches(2.55), Inches(2.9), Inches(0.3),
                label, font_size=11, bold=True, color=accent)
    add_textbox(s4, Inches(left.inches + 0.3), Inches(3.0), Inches(2.9), Inches(0.55),
                title, font_size=19, bold=True, color=TITLE_TEXT)
    add_textbox(s4, Inches(left.inches + 0.3), Inches(3.7), Inches(2.9), Inches(3.0),
                desc, font_size=14, color=SUB_TEXT, line_spacing=1.5)

# Arrows
for i in range(2):
    ax = Inches(4.55 + i * 3.85)
    add_textbox(s4, ax, Inches(4.2), Inches(0.4), Inches(0.4),
                ">>", font_size=22, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

add_slide_number(s4, 4)


# =============================================================================
#  SLIDE 5 - KEY FEATURES & IMPLEMENTATION  (White)
# =============================================================================
s5 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s5, SLIDE_BG)
add_thin_bar(s5, ACCENT_PURPLE)

section_tag(s5, Inches(1.0), Inches(0.65), "Key Features & Implementation", ACCENT_PURPLE)
slide_title(s5, Inches(1.0), Inches(1.1), "Core Functionalities in Plain Terms")

features = [
    ("Forward-Chaining\nInference Engine",
     "Systematically fires IF-THEN rules to walk a structured decision tree covering 9 fault categories.",
     ACCENT_CYAN),
    ("Interactive\nDiagnostic Console",
     "Split-pane UI: symptom input on the left, real-time system context on the right. Bento-box selectable options.",
     ACCENT_EMERALD),
    ("Live Telemetry\nDashboard",
     "Simulated real-time monitoring of CPU temp, Wi-Fi RSSI, heap memory, and GPIO states with oscilloscope waveforms.",
     ACCENT_AMBER),
    ("Searchable\nKnowledge Base",
     "9 categorized fault trees: Brownout, ESP-NOW, Wi-Fi, GPIO, Antenna, I2C, SPI, ADC, and Strapping Pin issues.",
     ACCENT_BLUE),
    ("ESP32 Pinout\nMapping",
     "Interactive 38-pin block diagram with real-time GPIO state overlay. Click any pin to inspect config and signal path.",
     ACCENT_PURPLE),
    ("Structured Diagnosis\nReports",
     "Each result includes root cause, confidence factor, severity rating, step-by-step fix, and circuit diagram notes.",
     ACCENT_ROSE),
]

for i, (title, desc, accent) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(1.0 + col * 3.85)
    top  = Inches(2.15 + row * 2.55)

    add_rect(s5, left, top, Inches(3.5), Inches(2.3), CARD_GRAY, CARD_BORDER)

    # Small accent bar at top of card
    bar = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.5), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    add_textbox(s5, Inches(left.inches + 0.25), Inches(top.inches + 0.2), Inches(3.0), Inches(0.65),
                title, font_size=14, bold=True, color=TITLE_TEXT, line_spacing=1.15)
    add_textbox(s5, Inches(left.inches + 0.25), Inches(top.inches + 1.0), Inches(3.0), Inches(1.15),
                desc, font_size=12, color=SUB_TEXT, line_spacing=1.4)

add_slide_number(s5, 5)


# =============================================================================
#  SLIDE 6 - RESULTS & DEMONSTRATION  (White)
# =============================================================================
s6 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s6, SLIDE_BG)
add_thin_bar(s6, ACCENT_EMERALD)

section_tag(s6, Inches(1.0), Inches(0.65), "Results & Demonstration", ACCENT_EMERALD)
slide_title(s6, Inches(1.0), Inches(1.1), "What Was Achieved & Tested")

# Metrics row
metrics = [
    ("2,400+", "Diagnostic Rules", ACCENT_CYAN),
    ("9", "Fault Categories", ACCENT_EMERALD),
    ("48+", "Hardware Profiles", ACCENT_AMBER),
    ("99.9%", "System Uptime", ACCENT_PURPLE),
]
for i, (val, lbl, clr) in enumerate(metrics):
    left = Inches(1.0 + i * 3.05)
    add_rect(s6, left, Inches(2.1), Inches(2.75), Inches(1.25), CARD_GRAY, CARD_BORDER)
    add_textbox(s6, Inches(left.inches + 0.1), Inches(2.2), Inches(2.55), Inches(0.55),
                val, font_size=28, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(s6, Inches(left.inches + 0.1), Inches(2.8), Inches(2.55), Inches(0.4),
                lbl, font_size=12, color=SUB_TEXT, align=PP_ALIGN.CENTER)

# Achievement bullets
add_rect(s6, Inches(1.0), Inches(3.65), Inches(11.3), Inches(3.35), CARD_GRAY, CARD_BORDER)

tf6 = add_rich_textbox(s6, Inches(1.5), Inches(3.85), Inches(10.3), Inches(3.0))

add_bullet(tf6,
    "Full Diagnostic Workflow Tested",
    "Users described symptoms and the system successfully navigated the decision tree to identify root causes across all 9 fault categories.",
    first=True)
add_bullet(tf6,
    "Professional SaaS-Grade UI",
    "Dark-mode dashboard with split-pane diagnostics, real-time waveforms, serial console, and ESP32 pinout diagram.")
add_bullet(tf6,
    "AI-Powered Chat Integration",
    "Integrated OpenRouter LLM API for natural-language follow-up questions alongside rule-based outputs.")
add_bullet(tf6,
    "Live Demo Available",
    "The application runs locally via Next.js dev server and is ready for live demonstration during the defense session.")

add_slide_number(s6, 6)


# =============================================================================
#  SLIDE 7 - PRACTICAL SIGNIFICANCE & IMPACT  (White)
# =============================================================================
s7 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s7, SLIDE_BG)
add_thin_bar(s7, ACCENT_AMBER)

section_tag(s7, Inches(1.0), Inches(0.65), "Practical Significance & Impact", ACCENT_AMBER)
slide_title(s7, Inches(1.0), Inches(1.1), "Why It Matters in Practice")

impacts = [
    ("Educational Value",
     "Gives ECE students a hands-on, interactive way to learn hardware debugging - bridging the gap between textbook theory and real-world embedded systems.",
     ACCENT_BLUE),
    ("Time Savings",
     "Reduces diagnostic time from hours of manual probing to minutes of guided, rule-based fault isolation - dramatically improving lab productivity.",
     ACCENT_EMERALD),
    ("Industry Relevance",
     "Mirrors professional IoT diagnostic platforms used in industry (AWS IoT, Vercel-style dashboards), preparing students for real engineering roles.",
     ACCENT_AMBER),
    ("Scalability & Extensibility",
     "The rule-based architecture can be extended to support additional microcontrollers (STM32, Arduino), new fault categories, and cloud-connected hardware.",
     ACCENT_PURPLE),
]

for i, (title, desc, accent) in enumerate(impacts):
    col = i % 2
    row = i // 2
    left = Inches(1.0 + col * 5.8)
    top  = Inches(2.15 + row * 2.55)

    add_rect(s7, left, top, Inches(5.5), Inches(2.3), CARD_GRAY, CARD_BORDER)

    # Accent bar
    bar = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), Inches(2.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    add_textbox(s7, Inches(left.inches + 0.35), Inches(top.inches + 0.25), Inches(4.8), Inches(0.4),
                title, font_size=18, bold=True, color=TITLE_TEXT)
    add_textbox(s7, Inches(left.inches + 0.35), Inches(top.inches + 0.8), Inches(4.8), Inches(1.3),
                desc, font_size=14, color=SUB_TEXT, line_spacing=1.5)

add_slide_number(s7, 7)


# =============================================================================
#  SLIDE 8 - GROUP MEMBERS  (White)
# =============================================================================
s8 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s8, SLIDE_BG)
add_thin_bar(s8, ACCENT_CYAN)

section_tag(s8, Inches(1.0), Inches(0.55), "Project Team  --  Group 11", ACCENT_CYAN)
slide_title(s8, Inches(1.0), Inches(0.95), "Meet the Team")

members = [
    ("Onyenaucheya Blessed Chimgozirim", "U2021/3020046", "Group Leader", "System Architecture & Forward-Chaining Inference Engine"),
    ("Memena Emmanuel Chiedu",           "U2021/3020054", "Member", "ESP-NOW Protocol & Peer MAC Fault Tree"),
    ("Ordu ThankGod Meyi",               "U2021/3020045", "Member", "Brownout & Power Supply Diagnostic Rules"),
    ("Paul Godwin",                       "U2021/3020047", "Member", "Wi-Fi Stack & FreeRTOS Watchdog Analysis"),
    ("Dickson Jessica Emem-Abasi",       "U2021/3020052", "Member", "GPIO Voltage Logic & Antenna Interference Domain"),
    ("Amadi Chibuike Eberechukwu",       "U2021/3020048", "Member", "I2C Bus Lockup & SPI Signal Integrity Testing"),
    ("Nwankwo Gift Chisom",               "U2021/3020049", "Member", "ADC2 Wi-Fi Conflict & Strapping Pin Validation"),
    ("Okonkwo Uchechukwu David",         "U2021/3020050", "Member", "Knowledge Base Data Entry & Rule Verification"),
    ("Justin Steve Homa",                "U2021/3020051", "Member", "Serial Monitor Panic Log Classifier Integration"),
    ("Okwudili Favour Chidinma",          "U2021/3020053", "Member", "UI/UX Design & Hardware Telemetry Dashboard"),
    ("Ugochukwu Emmanuel Kelechi",       "U2021/3020055", "Member", "Documentation, Testing & PDF Report Generation"),
]

# Header row
add_rect(s8, Inches(1.0), Inches(1.7), Inches(11.3), Inches(0.4), ACCENT_CYAN, None)
headers = [("S/N", 0.45), ("Full Name", 3.3), ("Matric No.", 1.6), ("Role", 1.2), ("Contribution", 4.2)]
x_pos = Inches(1.1)
for hdr, w in headers:
    add_textbox(s8, x_pos, Inches(1.72), Inches(w), Inches(0.35),
                hdr, font_size=11, bold=True, color=WHITE, font_name="Calibri")
    x_pos = Inches(x_pos.inches + w + 0.1)

# Member rows
for idx, (name, matric, role, contrib) in enumerate(members):
    row_top = Inches(2.2 + idx * 0.44)
    bg_color = CARD_GRAY if idx % 2 == 0 else SLIDE_BG

    add_rect(s8, Inches(1.0), row_top, Inches(11.3), Inches(0.42), bg_color, CARD_BORDER)

    # Highlight leader
    name_color = ACCENT_CYAN if role == "Group Leader" else TITLE_TEXT
    role_color  = ACCENT_CYAN if role == "Group Leader" else SUB_TEXT

    x = Inches(1.1)
    add_textbox(s8, x, Inches(row_top.inches + 0.03), Inches(0.45), Inches(0.35),
                str(idx + 1), font_size=10, bold=True, color=SUB_TEXT, align=PP_ALIGN.CENTER)
    x = Inches(x.inches + 0.55)
    add_textbox(s8, x, Inches(row_top.inches + 0.03), Inches(3.3), Inches(0.35),
                name, font_size=10, bold=(role == "Group Leader"), color=name_color)
    x = Inches(x.inches + 3.4)
    add_textbox(s8, x, Inches(row_top.inches + 0.03), Inches(1.6), Inches(0.35),
                matric, font_size=10, color=SUB_TEXT, font_name="Calibri")
    x = Inches(x.inches + 1.7)
    add_textbox(s8, x, Inches(row_top.inches + 0.03), Inches(1.2), Inches(0.35),
                role, font_size=10, bold=(role == "Group Leader"), color=role_color)
    x = Inches(x.inches + 1.3)
    add_textbox(s8, x, Inches(row_top.inches + 0.03), Inches(4.2), Inches(0.35),
                contrib, font_size=10, color=SUB_TEXT)

# Footer info
add_textbox(s8, Inches(1.0), Inches(7.0), Inches(11), Inches(0.3),
            "ECE 515.2 (Intro to AI)  |  Department of Electronic Engineering  |  University of Port Harcourt  |  2024/2025 Session",
            font_size=10, color=RGBColor(0xA0, 0xAE, 0xC0), align=PP_ALIGN.CENTER)

add_slide_number(s8, 8)


# =============================================================================
#  SLIDE 9 - CONCLUSION, FUTURE WORK & Q&A  (Dark navy)
# =============================================================================
s9 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s9, NAVY)

add_rect(s9, Inches(0), Inches(0), Inches(0.1), SLIDE_H, ACCENT_EMERALD, None)

section_tag(s9, Inches(1.2), Inches(0.7), "Conclusion  --  Future Work  --  Q&A", ACCENT_EMERALD)

add_textbox(s9, Inches(1.2), Inches(1.2), Inches(11), Inches(0.7),
            "Wrapping Up & Looking Ahead",
            font_size=34, bold=True, color=WHITE)

# Left - Conclusion
add_rect(s9, Inches(1.2), Inches(2.2), Inches(5.3), Inches(2.8), DARK_CARD, None)
add_textbox(s9, Inches(1.5), Inches(2.4), Inches(4.7), Inches(0.3),
            "CONCLUSION", font_size=12, bold=True, color=ACCENT_EMERALD)

tf9l = add_rich_textbox(s9, Inches(1.5), Inches(2.9), Inches(4.7), Inches(1.9))
add_bullet(tf9l,
    "Goal Achieved",
    "Successfully built a forward-chaining expert system for ESP32 hardware diagnostics.",
    font_size=14, lead_color=ACCENT_EMERALD, text_color=MEDIUM_GRAY, dot_color=ACCENT_EMERALD, first=True)
add_bullet(tf9l,
    "Classical AI Applied",
    "Demonstrated practical application of rule-based reasoning in embedded systems education.",
    font_size=14, lead_color=ACCENT_EMERALD, text_color=MEDIUM_GRAY, dot_color=ACCENT_EMERALD)

# Right - Future Work
add_rect(s9, Inches(6.9), Inches(2.2), Inches(5.3), Inches(2.8), DARK_CARD, None)
add_textbox(s9, Inches(7.2), Inches(2.4), Inches(4.7), Inches(0.3),
            "FUTURE WORK", font_size=12, bold=True, color=ACCENT_CYAN)

tf9r = add_rich_textbox(s9, Inches(7.2), Inches(2.9), Inches(4.7), Inches(1.9))
add_bullet(tf9r,
    "Hardware Integration",
    "Connect to real ESP32 boards via Web Serial API for live telemetry.",
    font_size=14, lead_color=ACCENT_CYAN, text_color=MEDIUM_GRAY, dot_color=ACCENT_CYAN, first=True)
add_bullet(tf9r,
    "Expand Rule Base",
    "Add support for STM32, Arduino, and Raspberry Pi with community fault trees.",
    font_size=14, lead_color=ACCENT_CYAN, text_color=MEDIUM_GRAY, dot_color=ACCENT_CYAN)
add_bullet(tf9r,
    "Cloud Deployment",
    "Deploy as a full SaaS platform with user accounts, saved diagnostics, and analytics.",
    font_size=14, lead_color=ACCENT_CYAN, text_color=MEDIUM_GRAY, dot_color=ACCENT_CYAN)

# Q&A
add_textbox(s9, Inches(1.2), Inches(5.5), Inches(11), Inches(0.7),
            "Questions & Answers",
            font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s9, Inches(1.2), Inches(6.2), Inches(11), Inches(0.5),
            "Thank you for your attention. We welcome your questions and feedback.",
            font_size=16, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

add_textbox(s9, Inches(1.2), Inches(6.9), Inches(11), Inches(0.4),
            "ECE 515.2  |  Group 11  |  University of Port Harcourt  |  2026",
            font_size=11, color=MUTED_GRAY, align=PP_ALIGN.CENTER)

add_slide_number(s9, 9, dark_bg=True)


# =============================================================================
#  SAVE
# =============================================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Final_Project_Defence.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
